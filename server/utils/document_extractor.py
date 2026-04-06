"""
utils/document_extractor.py

Unified text extractor for PDF and PPTX files.

PDF  → pymupdf4llm  — markdown-aware extraction (headings, tables, lists preserved)
                      Falls back to plain fitz if pymupdf4llm unavailable.
PPTX → python-pptx  — native slide-by-slide text extraction

Returns list of dicts: [{"page": int, "text": str}, ...]
so callers know which page/slide each chunk came from.

requirements.txt:
    pymupdf4llm>=0.0.17   (installs pymupdf automatically)
    python-pptx>=1.0.0
"""

import os
from logger import get_logger

logger = get_logger("document_extractor")


# ── PDF via pymupdf4llm ───────────────────────────────────────────────────────

def extract_pages_from_pdf(path: str) -> list[dict]:
    """
    Extract text page-by-page from a PDF using pymupdf4llm.

    pymupdf4llm.to_markdown(path, page_chunks=True) returns:
        [
            {
                "metadata": {"page": 1, ...},
                "text": "# Heading\n\nParagraph text...",
                ...
            },
            ...
        ]

    The markdown output preserves headings, bullet lists, and table structure,
    which improves chunking quality for RAG pipelines vs raw plain text.

    Falls back to plain fitz extraction if pymupdf4llm is not installed.
    """
    try:
        import pymupdf4llm
    except ImportError:
        logger.warning("pymupdf4llm not found, falling back to plain fitz. "
                       "Add 'pymupdf4llm' to requirements.txt for better extraction.")
        return _extract_pages_fitz_fallback(path)

    try:
        page_chunks = pymupdf4llm.to_markdown(path, page_chunks=True)
    except Exception as e:
        logger.warning(f"pymupdf4llm failed ({e}), falling back to plain fitz.")
        return _extract_pages_fitz_fallback(path)

    pages = []
    for chunk in page_chunks:
        text = (chunk.get("text") or "").strip()
        if not text:
            continue
        # page_chunks metadata uses 0-based page numbers; normalise to 1-based
        raw_page = (chunk.get("metadata") or {}).get("page", 0)
        pages.append({
            "page": int(raw_page) + 1 if isinstance(raw_page, int) else int(raw_page or 1),
            "text": text
        })

    total_chars = sum(len(p["text"]) for p in pages)
    logger.info(f"PDF extracted (pymupdf4llm): {len(pages)} pages, "
                f"{total_chars} chars from {os.path.basename(path)}")
    return pages


def _extract_pages_fitz_fallback(path: str) -> list[dict]:
    """Plain fitz fallback — used only if pymupdf4llm is unavailable."""
    try:
        import fitz
    except ImportError:
        raise RuntimeError(
            "Neither pymupdf4llm nor pymupdf (fitz) is installed. "
            "Add 'pymupdf4llm' to requirements.txt"
        )

    pages = []
    doc = fitz.open(path)
    for i, page in enumerate(doc):
        text = page.get_text("text")
        if text and text.strip():
            pages.append({"page": i + 1, "text": text.strip()})
    doc.close()

    total_chars = sum(len(p["text"]) for p in pages)
    logger.info(f"PDF extracted (fitz fallback): {len(pages)} pages, "
                f"{total_chars} chars from {os.path.basename(path)}")
    return pages


# ── PPTX via python-pptx ──────────────────────────────────────────────────────

def extract_pages_from_pptx(path: str) -> list[dict]:
    """
    Extract text slide-by-slide from a PPTX file.
    Returns [{"page": 1, "text": "..."}, ...]  (page = slide number)
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        raise RuntimeError("python-pptx is not installed. Add 'python-pptx' to requirements.txt")

    prs = Presentation(path)
    slides = []

    for i, slide in enumerate(prs.slides):
        parts = []

        # Slide title first (if present)
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(f"[Title] {slide.shapes.title.text.strip()}")

        # All text frames in reading order
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Skip title — already handled above
            if shape == slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    # Preserve bullet hierarchy with indentation
                    indent = "  " * (para.level or 0)
                    parts.append(f"{indent}{line}")

        # Speaker notes (useful context for RAG)
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Notes] {notes_text}")

        text = "\n".join(parts)
        if text.strip():
            slides.append({
                "page": i + 1,
                "text": text.strip()
            })

    total_chars = sum(len(s["text"]) for s in slides)
    logger.info(f"PPTX extracted: {len(slides)} slides, {total_chars} chars from {os.path.basename(path)}")

    return slides


# ── Unified entry point ───────────────────────────────────────────────────────

def extract_pages(path: str) -> list[dict]:
    """
    Auto-detect file type and extract text.
    Returns [{"page": int, "text": str}, ...]
    """
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return extract_pages_from_pdf(path)
    elif ext in (".pptx",):
        return extract_pages_from_pptx(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: .pdf, .pptx")


def extract_text(path: str) -> str:
    """
    Convenience wrapper — returns full text as single string.
    Used by callers that don't need per-page info (e.g. topic extraction).
    """
    pages = extract_pages(path)
    return "\n\n".join(
        f"Page {p['page']}:\n{p['text']}" for p in pages
    )


def get_supported_extensions() -> list[str]:
    return [".pdf", ".pptx"]